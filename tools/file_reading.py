# core/file_reading.py - File attachment processing untuk Mirai
"""
Module untuk membaca dan mengekstrak teks dari berbagai format file.
Mendukung: PDF, DOCX, XLSX, PPTX, TXT
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation
from config import (
    SUPPORTED_EXTENSIONS, MAX_ATTACHMENTS, MAX_FILE_SIZE_BYTES,
    MAX_TEXT_PER_FILE_CHARS, MAX_TOTAL_CHARS
)

def _clip(text: str, limit: int) -> tuple[str, bool]:
    """
    Potong teks jika melebihi limit.
    
    Args:
        text: Teks yang akan dipotong
        limit: Batas karakter
        
    Returns:
        Tuple of (clipped_text, was_clipped)
    """
    if len(text) <= limit:
        return text, False
    return text[:limit].rstrip() + "\n...[truncated]", True


def _extract_pdf(raw: bytes) -> str:
    """Extract teks dari PDF file."""
    reader = PdfReader(BytesIO(raw))
    parts: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            parts.append(page_text.strip())
    return "\n\n".join(parts)


def _extract_docx(raw: bytes) -> str:
    """Extract teks dari DOCX file."""
    doc = Document(BytesIO(raw))
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(lines)


def _extract_xlsx(raw: bytes) -> str:
    """Extract teks dari XLSX file."""
    wb = load_workbook(filename=BytesIO(raw), data_only=True, read_only=True)
    chunks: list[str] = []
    for sheet in wb.worksheets:
        chunks.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if row_values:
                chunks.append(" | ".join(row_values))
        chunks.append("")
    return "\n".join(chunks).strip()


def _extract_pptx(raw: bytes) -> str:
    """Extract teks dari PPTX file."""
    prs = Presentation(BytesIO(raw))
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_lines.append(text.strip())
        if slide_lines:
            chunks.append(f"[Slide {idx}]")
            chunks.extend(slide_lines)
            chunks.append("")
    return "\n".join(chunks).strip()


def _extract_txt(raw: bytes) -> str:
    """Extract teks dari plain text file."""
    try:
        # Coba decode sebagai UTF-8 (default)
        return raw.decode('utf-8').strip()
    except UnicodeDecodeError:
        # Fallback ke latin-1 jika UTF-8 gagal
        try:
            return raw.decode('latin-1').strip()
        except:
            # Jika semua gagal, return string dari bytes dengan ignore error
            return raw.decode('utf-8', errors='ignore').strip()


def extract_file_text(filename: str, raw: bytes) -> str:
    """
    Extract teks dari file berdasarkan ekstensi.
    
    Args:
        filename: Nama file
        raw: Raw bytes dari file
        
    Returns:
        str: Extracted text
    """
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(raw)
    if ext == ".docx":
        return _extract_docx(raw)
    if ext == ".xlsx":
        return _extract_xlsx(raw)
    if ext == ".pptx":
        return _extract_pptx(raw)
    if ext == ".txt":
        return _extract_txt(raw)
    return ""


async def build_attachment_context(attachments: Iterable) -> str:
    """
    Build konteks dari attachment files.
    
    Args:
        attachments: Iterable of Discord attachment objects
        
    Returns:
        str: Formatted attachment context untuk dikirim ke AI
    """
    docs = list(attachments)[:MAX_ATTACHMENTS]
    if not docs:
        return ""

    collected: list[str] = []
    notes: list[str] = []
    total_chars = 0

    for attachment in docs:
        filename = getattr(attachment, "filename", "unknown")
        size = int(getattr(attachment, "size", 0) or 0)
        ext = Path(filename).suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            notes.append(f"- `{filename}` dilewati (format tidak didukung).")
            continue

        if size > MAX_FILE_SIZE_BYTES:
            notes.append(f"- `{filename}` dilewati (ukuran > 10MB).")
            continue

        try:
            raw = await attachment.read()
            extracted = extract_file_text(filename, raw).strip()
            if not extracted:
                notes.append(f"- `{filename}` tidak punya teks yang bisa diekstrak.")
                continue

            extracted, _ = _clip(extracted, MAX_TEXT_PER_FILE_CHARS)
            remaining = MAX_TOTAL_CHARS - total_chars
            if remaining <= 0:
                notes.append("- Batas total teks attachment tercapai, file berikutnya dilewati.")
                break

            extracted, total_clipped = _clip(extracted, remaining)
            if total_clipped:
                notes.append("- Teks attachment dipotong karena melebihi batas total.")

            total_chars += len(extracted)
            collected.append(f"[File: {filename}]\n{extracted}")
        except Exception as err:
            notes.append(f"- `{filename}` gagal dibaca ({err}).")

    if not collected and not notes:
        return ""

    sections: list[str] = ["[Attachment Context]"]
    if collected:
        sections.append("\n\n".join(collected))
    if notes:
        sections.append("Catatan:\n" + "\n".join(notes))
    return "\n\n".join(sections)
